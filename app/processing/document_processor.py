import io
import os
import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import tempfile
import mimetypes

# 텍스트 추출 라이브러리들
try:
    import PyPDF2
    from pdfplumber import PDF
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """문서 처리 및 텍스트 추출 클래스"""
    
    SUPPORTED_FORMATS = {
        'pdf': ['application/pdf', '.pdf'],
        'docx': ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', '.docx'],
        'doc': ['application/msword', '.doc'],
        'pptx': ['application/vnd.openxmlformats-officedocument.presentationml.presentation', '.pptx'],
        'ppt': ['application/vnd.ms-powerpoint', '.ppt'],
        'txt': ['text/plain', '.txt'],
        'xlsx': ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', '.xlsx'],
        'csv': ['text/csv', '.csv'],
        'md': ['text/markdown', '.md'],
        'html': ['text/html', '.html']
    }
    
    def __init__(self, upload_dir: str = "app/uploads"):
        """
        초기화
        Args:
            upload_dir: 업로드 파일 저장 디렉토리
        """
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        
        # 임시 파일 디렉토리
        self.temp_dir = Path(tempfile.gettempdir()) / "education_platform"
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"문서 처리기 초기화 완료 - 업로드 디렉토리: {self.upload_dir}")
    
    def get_supported_formats(self) -> List[str]:
        """지원되는 파일 형식 목록 반환"""
        return list(self.SUPPORTED_FORMATS.keys())
    
    def is_supported_format(self, file_path: str) -> bool:
        """파일 형식 지원 여부 확인"""
        file_ext = Path(file_path).suffix.lower()
        for format_info in self.SUPPORTED_FORMATS.values():
            if file_ext in format_info:
                return True
        return False
    
    def detect_file_type(self, file_path: str) -> str:
        """파일 타입 감지"""
        file_ext = Path(file_path).suffix.lower()
        
        for file_type, format_info in self.SUPPORTED_FORMATS.items():
            if file_ext in format_info:
                return file_type
        
        return 'unknown'
    
    def save_uploaded_file(self, uploaded_file, course_id: str, user_id: str) -> Tuple[str, Dict]:
        """
        업로드된 파일 저장
        Args:
            uploaded_file: Streamlit 업로드 파일 객체
            course_id: 강의 ID
            user_id: 사용자 ID
        Returns:
            (파일 경로, 메타데이터)
        """
        try:
            # 파일 저장 경로 생성
            course_dir = self.upload_dir / course_id
            course_dir.mkdir(parents=True, exist_ok=True)
            
            # 파일명 생성 (중복 방지)
            file_extension = Path(uploaded_file.name).suffix
            safe_filename = f"{user_id}_{uploaded_file.name}"
            file_path = course_dir / safe_filename
            
            # 중복 파일명 처리
            counter = 1
            while file_path.exists():
                name_without_ext = Path(uploaded_file.name).stem
                safe_filename = f"{user_id}_{name_without_ext}_{counter}{file_extension}"
                file_path = course_dir / safe_filename
                counter += 1
            
            # 파일 저장
            with open(file_path, 'wb') as f:
                f.write(uploaded_file.getvalue())
            
            # 메타데이터 생성
            metadata = {
                'original_filename': uploaded_file.name,
                'saved_filename': safe_filename,
                'file_path': str(file_path),
                'file_size': uploaded_file.size,
                'file_type': uploaded_file.type,
                'course_id': course_id,
                'uploaded_by': user_id,
                'detected_type': self.detect_file_type(str(file_path))
            }
            
            logger.info(f"파일 저장 완료: {file_path}")
            return str(file_path), metadata
            
        except Exception as e:
            logger.error(f"파일 저장 중 오류 발생: {str(e)}")
            raise
    
    def extract_text_from_file(self, file_path: str) -> Dict:
        """
        파일에서 텍스트 추출
        Args:
            file_path: 파일 경로
        Returns:
            추출 결과 딕셔너리
        """
        try:
            file_type = self.detect_file_type(file_path)
            
            if file_type == 'pdf':
                return self._extract_from_pdf(file_path)
            elif file_type == 'docx':
                return self._extract_from_docx(file_path)
            elif file_type == 'pptx':
                return self._extract_from_pptx(file_path)
            elif file_type == 'txt':
                return self._extract_from_txt(file_path)
            elif file_type == 'xlsx':
                return self._extract_from_xlsx(file_path)
            elif file_type == 'csv':
                return self._extract_from_csv(file_path)
            elif file_type in ['md', 'html']:
                return self._extract_from_txt(file_path)
            else:
                return {
                    'success': False,
                    'text': '',
                    'error': f'지원되지 않는 파일 형식: {file_type}',
                    'page_count': 0,
                    'word_count': 0
                }
                
        except Exception as e:
            logger.error(f"텍스트 추출 중 오류 발생: {str(e)}")
            return {
                'success': False,
                'text': '',
                'error': str(e),
                'page_count': 0,
                'word_count': 0
            }
    
    def _extract_from_pdf(self, file_path: str) -> Dict:
        """PDF 파일에서 텍스트 추출"""
        if not PYPDF2_AVAILABLE:
            return {
                'success': False,
                'text': '',
                'error': 'PDF 처리 라이브러리가 설치되지 않았습니다.',
                'page_count': 0,
                'word_count': 0
            }
        
        try:
            text = ""
            page_count = 0
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                page_count = len(pdf_reader.pages)
                
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            
            # 텍스트 정리
            text = self._clean_text(text)
            word_count = len(text.split())
            
            return {
                'success': True,
                'text': text,
                'error': None,
                'page_count': page_count,
                'word_count': word_count
            }
            
        except Exception as e:
            logger.error(f"PDF 처리 중 오류: {str(e)}")
            return {
                'success': False,
                'text': '',
                'error': str(e),
                'page_count': 0,
                'word_count': 0
            }
    
    def _extract_from_docx(self, file_path: str) -> Dict:
        """DOCX 파일에서 텍스트 추출"""
        if not DOCX_AVAILABLE:
            return {
                'success': False,
                'text': '',
                'error': 'DOCX 처리 라이브러리가 설치되지 않았습니다.',
                'page_count': 0,
                'word_count': 0
            }
        
        try:
            doc = Document(file_path)
            text = ""
            
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # 테이블 내용도 추출
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            text = self._clean_text(text)
            word_count = len(text.split())
            
            return {
                'success': True,
                'text': text,
                'error': None,
                'page_count': 1,  # DOCX는 페이지 수 계산이 복잡
                'word_count': word_count
            }
            
        except Exception as e:
            logger.error(f"DOCX 처리 중 오류: {str(e)}")
            return {
                'success': False,
                'text': '',
                'error': str(e),
                'page_count': 0,
                'word_count': 0
            }
    
    def _extract_from_pptx(self, file_path: str) -> Dict:
        """PPTX 파일에서 텍스트 추출"""
        if not PPTX_AVAILABLE:
            return {
                'success': False,
                'text': '',
                'error': 'PPTX 처리 라이브러리가 설치되지 않았습니다.',
                'page_count': 0,
                'word_count': 0
            }
        
        try:
            prs = Presentation(file_path)
            text = ""
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            text = self._clean_text(text)
            word_count = len(text.split())
            
            return {
                'success': True,
                'text': text,
                'error': None,
                'page_count': slide_count,
                'word_count': word_count
            }
            
        except Exception as e:
            logger.error(f"PPTX 처리 중 오류: {str(e)}")
            return {
                'success': False,
                'text': '',
                'error': str(e),
                'page_count': 0,
                'word_count': 0
            }
    
    def _extract_from_txt(self, file_path: str) -> Dict:
        """텍스트 파일에서 내용 추출"""
        try:
            encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']
            text = ""
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if not text:
                raise Exception("텍스트 파일 인코딩을 인식할 수 없습니다.")
            
            text = self._clean_text(text)
            word_count = len(text.split())
            
            return {
                'success': True,
                'text': text,
                'error': None,
                'page_count': 1,
                'word_count': word_count
            }
            
        except Exception as e:
            logger.error(f"텍스트 파일 처리 중 오류: {str(e)}")
            return {
                'success': False,
                'text': '',
                'error': str(e),
                'page_count': 0,
                'word_count': 0
            }
    
    def _extract_from_xlsx(self, file_path: str) -> Dict:
        """XLSX 파일에서 텍스트 추출"""
        if not OPENPYXL_AVAILABLE:
            return {
                'success': False,
                'text': '',
                'error': 'XLSX 처리 라이브러리가 설치되지 않았습니다.',
                'page_count': 0,
                'word_count': 0
            }
        
        try:
            workbook = openpyxl.load_workbook(file_path)
            text = ""
            sheet_count = 0
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_count += 1
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        text += row_text + "\n"
            
            text = self._clean_text(text)
            word_count = len(text.split())
            
            return {
                'success': True,
                'text': text,
                'error': None,
                'page_count': sheet_count,
                'word_count': word_count
            }
            
        except Exception as e:
            logger.error(f"XLSX 처리 중 오류: {str(e)}")
            return {
                'success': False,
                'text': '',
                'error': str(e),
                'page_count': 0,
                'word_count': 0
            }
    
    def _extract_from_csv(self, file_path: str) -> Dict:
        """CSV 파일에서 텍스트 추출"""
        try:
            import csv
            
            text = ""
            row_count = 0
            
            encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding, newline='') as file:
                        csv_reader = csv.reader(file)
                        for row in csv_reader:
                            row_count += 1
                            text += " ".join(row) + "\n"
                    break
                except UnicodeDecodeError:
                    continue
            
            if not text:
                raise Exception("CSV 파일 인코딩을 인식할 수 없습니다.")
            
            text = self._clean_text(text)
            word_count = len(text.split())
            
            return {
                'success': True,
                'text': text,
                'error': None,
                'page_count': 1,
                'word_count': word_count
            }
            
        except Exception as e:
            logger.error(f"CSV 처리 중 오류: {str(e)}")
            return {
                'success': False,
                'text': '',
                'error': str(e),
                'page_count': 0,
                'word_count': 0
            }
    
    def _clean_text(self, text: str) -> str:
        """텍스트 정리"""
        if not text:
            return ""
        
        # 연속된 공백 정리
        import re
        text = re.sub(r'\s+', ' ', text)
        
        # 연속된 줄바꿈 정리
        text = re.sub(r'\n+', '\n', text)
        
        # 앞뒤 공백 제거
        text = text.strip()
        
        return text
    
    def get_file_info(self, file_path: str) -> Dict:
        """파일 정보 조회"""
        try:
            file_path = Path(file_path)
            
            if not file_path.exists():
                return {'error': '파일을 찾을 수 없습니다.'}
            
            stat = file_path.stat()
            
            return {
                'filename': file_path.name,
                'size': stat.st_size,
                'size_mb': stat.st_size / (1024 * 1024),
                'modified': stat.st_mtime,
                'file_type': self.detect_file_type(str(file_path)),
                'is_supported': self.is_supported_format(str(file_path))
            }
            
        except Exception as e:
            logger.error(f"파일 정보 조회 중 오류: {str(e)}")
            return {'error': str(e)}
    
    def delete_file(self, file_path: str) -> bool:
        """파일 삭제"""
        try:
            file_path = Path(file_path)
            if file_path.exists():
                file_path.unlink()
                logger.info(f"파일 삭제 완료: {file_path}")
                return True
            return False
            
        except Exception as e:
            logger.error(f"파일 삭제 중 오류: {str(e)}")
            return False 